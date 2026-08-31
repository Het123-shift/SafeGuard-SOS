import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- COLOR PALETTE ---
COLOR_BG = RGBColor(9, 13, 22)           # Dark Navy/Slate (#090D16)
COLOR_CARD_BG = RGBColor(22, 29, 47)     # Glass Card BG (#161D2F)
COLOR_CARD_BORDER = RGBColor(48, 54, 61) # Border (#30363D)
COLOR_RED = RGBColor(255, 77, 77)        # Emergency Crimson (#FF4D4D)
COLOR_GREEN = RGBColor(16, 185, 129)     # Emerald Green (#10B981)
COLOR_CYAN = RGBColor(56, 189, 248)      # Cyber Sky / Cyan (#38BDF8)
COLOR_AMBER = RGBColor(251, 191, 36)     # Warning Amber (#FBBF24)
COLOR_PURPLE = RGBColor(167, 139, 250)   # Purple (#A78BFA)
COLOR_WHITE = RGBColor(255, 255, 255)    # White
COLOR_MUTED = RGBColor(148, 163, 184)    # Slate Muted (#94A3B8)
COLOR_CARD_RED = RGBColor(45, 15, 15)    # Dark Red Card (#2D0F0F)
COLOR_CARD_GREEN = RGBColor(10, 35, 20)  # Dark Green Card (#0A2314)
COLOR_CARD_AMBER = RGBColor(40, 28, 8)   # Dark Amber Card (#281C08)

BASE_DIR = os.getcwd()
DASHBOARD_IMG = os.path.join(BASE_DIR, "dashboard_screenshot.png")
SCREEN_IMG = os.path.join(BASE_DIR, "screen.png")
LOGO_IMG = os.path.join(BASE_DIR, "assets", "images", "logo.png")

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.color.rgb = COLOR_BG
        return bg

    def add_header(slide, title_text, category_text="TEAM PIXEL PERFECT (SM083)  |  HACKATHON ROUND 2"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_CARD_BORDER
        line.line.color.rgb = COLOR_CARD_BORDER

    def add_card(slide, left, top, width, height, title, subtitle, bullets, card_bg=COLOR_CARD_BG, accent_color=COLOR_CYAN):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1)

        padding = Inches(0.2)
        tb = slide.shapes.add_textbox(left + padding, top + padding, width - (padding * 2), height - (padding * 2))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = accent_color
        p0.space_after = Pt(4)

        if subtitle:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.size = Pt(11)
            p_sub.font.italic = True
            p_sub.font.color.rgb = COLOR_MUTED
            p_sub.space_after = Pt(6)

        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = "•  " + b
            pb.font.size = Pt(11.5)
            pb.font.color.rgb = COLOR_WHITE
            pb.space_after = Pt(3)

    # -------------------------------------------------------------
    # SLIDE 1: WHAT WAS OUR IDEATION (WITH LOGO)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    glow = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    glow.fill.solid()
    glow.fill.fore_color.rgb = COLOR_CARD_BG
    glow.line.color.rgb = COLOR_RED
    glow.line.width = Pt(2)

    # Add Logo if present
    if os.path.exists(LOGO_IMG):
        try:
            s1.shapes.add_picture(LOGO_IMG, Inches(10.5), Inches(1.1), Inches(1.6), Inches(1.6))
        except Exception as e:
            print("Logo insert error:", e)

    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(9.2), Inches(5.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p_badge = tf1.paragraphs[0]
    p_badge.text = "🚨 HACKATHON ROUND 2 FINALIST PITCH  |  TRACK: SMART SAFETY & IOT"
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_RED
    p_badge.space_after = Pt(6)

    p_main = tf1.add_paragraph()
    p_main.text = "SafeGuard SOS"
    p_main.font.size = Pt(42)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_WHITE
    p_main.space_after = Pt(4)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Autonomous Emergency Dispatch & Edge Telemetry Ecosystem"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_CYAN
    p_sub.space_after = Pt(10)

    p_ideation = tf1.add_paragraph()
    p_ideation.text = "Our Ideation: Traditional safety apps fail because they assume ideal conditions (5G data, unlocked screens, calm victims). We envisioned an autonomous, zero-friction safety pipeline that triggers without screen interaction, requires ZERO app install for responders, and delivers tamper-proof forensic evidence."
    p_ideation.font.size = Pt(13)
    p_ideation.font.color.rgb = COLOR_MUTED
    p_ideation.space_after = Pt(16)

    # TEAM INFO BOX ON SLIDE 1
    p_team_header = tf1.add_paragraph()
    p_team_header.text = "👥 TEAM: PIXEL PERFECT (TEAM ID: SM083)"
    p_team_header.font.size = Pt(15)
    p_team_header.font.bold = True
    p_team_header.font.color.rgb = COLOR_GREEN
    p_team_header.space_after = Pt(8)

    members = [
        ("HET PATEL", "240410149035"),
        ("RUDRA PATEL", "240410149043"),
        ("PRITAM PATEL", "240410149047"),
        ("VEDANSHU PATEL", "240410149044")
    ]
    
    mem_text = "   •   ".join([f"{name} ({id_num})" for name, id_num in members])
    p_mem = tf1.add_paragraph()
    p_mem.text = mem_text
    p_mem.font.size = Pt(12)
    p_mem.font.bold = True
    p_mem.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 2: THE CRITICAL PROBLEM
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Slide 02: The Critical Problem — Why Safety Apps Fail Under Duress")

    add_card(s2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "1. High Panic Latency",
             "Cognitive overload during attacks/accidents",
             [
                 "Victims cannot unlock screens or dial numbers under physical assault.",
                 "Over 70% of victims are physically unable to articulate GPS coordinates.",
                 "Single-button apps still require unlocking the OS first."
             ],
             accent_color=COLOR_RED)

    add_card(s2, Inches(4.85), Inches(1.8), Inches(3.6), Inches(5.0),
             "2. The Connectivity Trap",
             "Single-point-of-failure infrastructure",
             [
                 "Most apps fail completely when 4G/5G data drops in basements or transit.",
                 "Cloud-only dispatch stalls when backend webhooks timeout.",
                 "WhatsApp/VoIP deep links fail without active internet."
             ],
             accent_color=COLOR_AMBER)

    add_card(s2, Inches(8.9), Inches(1.8), Inches(3.6), Inches(5.0),
             "3. Tampered Evidence",
             "Zero forensic custody during crimes",
             [
                 "Perpetrators often smash or discard victim phones immediately.",
                 "Audio and photo evidence captured unencrypted is lost or inaccessible.",
                 "Families and guardians have no immediate access to locked incident evidence."
             ],
             accent_color=COLOR_PURPLE)

    # -------------------------------------------------------------
    # SLIDE 3: THE SAFEGUARD SOLUTION
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "Slide 03: The SafeGuard Solution — Multi-Vector Resilient Ecosystem")

    add_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4),
             "Multi-Modal Emergency Triggers",
             "Trigger anywhere, anytime without friction",
             [
                 "One-touch instant SOS button with 3-second abort timer.",
                 "Continuous accelerometer-based fall & crash impact detection.",
                 "Hardware volume key sequences & Smartwatch BLE wrist trigger."
             ],
             accent_color=COLOR_RED)

    add_card(s3, Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.4),
             "Resilient Multi-Channel Dispatch",
             "Zero Single Point of Failure (SPOF)",
             [
                 "Direct background SIM SMS sent over cellular bands (No internet needed).",
                 "International E.164 WhatsApp auto deep-link bridge.",
                 "Automated siren activation + silent dispatch modes."
             ],
             accent_color=COLOR_GREEN)

    add_card(s3, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4),
             "Zero-Install Responder Web Portal",
             "Instant situational awareness for contacts & police",
             [
                 "Responders open an ephemeral live map on any mobile browser.",
                 "No app installation required for parents or emergency responders.",
                 "1-Click routing directly into Google Maps / Apple Maps."
             ],
             accent_color=COLOR_CYAN)

    add_card(s3, Inches(6.9), Inches(4.5), Inches(5.6), Inches(2.4),
             "AES-256 Encrypted Evidence Vault",
             "Tamper-proof black-box incident recording",
             [
                 "Automatic ambient microphone audio recording on trigger.",
                 "Encrypted local storage with PIN / Biometric security lock.",
                 "Emergency Medical ID broadcast (Blood group, allergies, conditions)."
             ],
             accent_color=COLOR_PURPLE)

    # -------------------------------------------------------------
    # SLIDE 4: WHAT MVP WE ACHIEVED (WITH APP SCREENSHOT)
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Slide 04: What MVP We Achieved — Operational Round 2 Capabilities")

    add_card(s4, Inches(0.8), Inches(1.8), Inches(4.2), Inches(5.0),
             "Active SOS & Telephony Engine",
             "Verified on Android hardware & Web",
             [
                 "3s countdown with abort modal & haptics.",
                 "Audible pulse siren with instant silent toggle.",
                 "Native Android SEND_SMS direct SIM dispatch.",
                 "WhatsApp wa.me link with E.164 normalization.",
                 "Medical ID transmission (Blood, allergies)."
             ],
             accent_color=COLOR_GREEN)

    add_card(s4, Inches(5.2), Inches(1.8), Inches(4.2), Inches(5.0),
             "Evidence Vault & Wearables",
             "Secure evidence & peripheral layer",
             [
                 "AES-256 encrypted vault for audio & media.",
                 "Ambient audio recording during active SOS.",
                 "Smartwatch BLE scanning & discovery.",
                 "Family Circle real-time status dashboard."
             ],
             accent_color=COLOR_CYAN)

    # Insert screenshot card on right side
    if os.path.exists(DASHBOARD_IMG):
        try:
            # Card background for screenshot
            card_img = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.6), Inches(1.8), Inches(2.933), Inches(5.0))
            card_img.fill.solid()
            card_img.fill.fore_color.rgb = COLOR_CARD_BG
            card_img.line.color.rgb = COLOR_CYAN
            card_img.line.width = Pt(1.5)
            s4.shapes.add_picture(DASHBOARD_IMG, Inches(9.75), Inches(2.0), Inches(2.633), Inches(4.6))
        except Exception as e:
            print("Dashboard image insert error:", e)

    # -------------------------------------------------------------
    # SLIDE 5: TRANSPARENT ENGINEERING ANALYSIS
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Slide 05: Transparent Engineering Analysis — Intentional Fallback Matrix")

    add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4),
             "1. Geolocation: Last-Known Location Fallback",
             "Prevents satellite GPS stall delays",
             [
                 "Current Behavior: Transmits immediate verified Last Known Coordinates.",
                 "Why it's a Fallback: Prevents 15-second satellite acquisition stalls in basements/elevators so the SOS fires without delay."
             ],
             card_bg=COLOR_CARD_AMBER, accent_color=COLOR_AMBER)

    add_card(s5, Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.4),
             "2. Connectivity: Native SIM SMS Fallback",
             "Graceful degradation when data is down",
             [
                 "Current Behavior: WhatsApp deep-links require internet.",
                 "Why it's a Fallback: If offline, SafeGuard automatically falls back to direct native cellular SIM SMS."
             ],
             card_bg=COLOR_CARD_GREEN, accent_color=COLOR_GREEN)

    add_card(s5, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4),
             "3. Evidence Vault: Master Passkey Access",
             "Cryptographic local protection",
             [
                 "Current Behavior: Audio files are strongly encrypted in the vault.",
                 "Why it's a Fallback: Requires family/guardian to enter passkey to prevent attacker tampering at the scene."
             ],
             card_bg=COLOR_CARD_BG, accent_color=COLOR_CYAN)

    add_card(s5, Inches(6.9), Inches(4.5), Inches(5.6), Inches(2.4),
             "4. Sensor & Hardware Baseline Thresholds",
             "Conservative calibration to avoid false alarms",
             [
                 "Current Behavior: High G-force impact threshold & in-app volume 3-tap.",
                 "Why it's a Fallback: Eliminates accidental false dispatches during normal running or minor phone handling."
             ],
             card_bg=COLOR_CARD_BG, accent_color=COLOR_MUTED)

    # -------------------------------------------------------------
    # SLIDE 6: SYSTEM ARCHITECTURE
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Slide 06: System Architecture — Modular, Reactive & Edge-Optimized")

    add_card(s6, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Mobile & Web Client",
             "Cross-Platform Frontend",
             [
                 "React Native 0.79 & Expo SDK 53.",
                 "TypeScript for strict type safety.",
                 "Expo Router 5 dynamic routing.",
                 "Expo Sensors (Accelerometer @ 100ms polling interval)."
             ],
             accent_color=COLOR_CYAN)

    add_card(s6, Inches(4.85), Inches(1.8), Inches(3.6), Inches(5.0),
             "Backend & Realtime Edge",
             "Supabase Cloud Infrastructure",
             [
                 "PostgreSQL database with Row-Level Security (RLS).",
                 "Realtime WebSocket channels for live geolocation broadcasting.",
                 "Encrypted buckets for evidence backup.",
                 "Serverless Edge Functions for multi-channel SMS webhooks."
             ],
             accent_color=COLOR_GREEN)

    add_card(s6, Inches(8.9), Inches(1.8), Inches(3.6), Inches(5.0),
             "Native Hardware Bridges",
             "Android OS & Sensor Hooks",
             [
                 "Android Telephony SMSManager for direct SIM SMS.",
                 "Web Audio API & Native Audio Recorder for ambient mic capture.",
                 "BLE Watch Connectivity layer for smartwatch sync.",
                 "Local SQLite / AsyncStorage for offline incident caching."
             ],
             accent_color=COLOR_PURPLE)

    # -------------------------------------------------------------
    # SLIDE 7: WHERE IT CHOKES
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Slide 07: Where It Chokes — Honest Stress-Test Bottlenecks & Gaps", "ENGINEERING CHOKE-POINTS")

    add_card(s7, Inches(0.8), Inches(1.8), Inches(3.6), Inches(2.4),
             "1. Cold-Start Latency",
             "Bundle initialization delay",
             [
                 "Takes too much time to initialize JS bundle on sudden launch under duress."
             ],
             card_bg=COLOR_CARD_RED, accent_color=COLOR_RED)

    add_card(s7, Inches(4.85), Inches(1.8), Inches(3.6), Inches(2.4),
             "2. Static Geolocation",
             "Telemetry continuity gap",
             [
                 "Transmits static Last Known Location rather than continuous live dynamic stream."
             ],
             card_bg=COLOR_CARD_AMBER, accent_color=COLOR_AMBER)

    add_card(s7, Inches(8.9), Inches(1.8), Inches(3.6), Inches(2.4),
             "3. Offline WhatsApp Limit",
             "Internet dependency",
             [
                 "WhatsApp deep-links fail without data; relies solely on single-channel SIM SMS."
             ],
             card_bg=COLOR_CARD_BG, accent_color=COLOR_PURPLE)

    add_card(s7, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4),
             "4. Evidence Vault Decryption Barrier",
             "Local key access constraint",
             [
                 "Audio files are heavily encrypted locally; guardians cannot access recordings remotely without physical access and knowing the victim's device passkey."
             ],
             card_bg=COLOR_CARD_BG, accent_color=COLOR_CYAN)

    add_card(s7, Inches(6.9), Inches(4.5), Inches(5.6), Inches(2.4),
             "5. OS Lock-Screen Hook & High Fall Force",
             "OS background event & sensor calibration limits",
             [
                 "Android OS blocks 3-tap volume key events when phone screen is locked; fall detection sensitivity threshold is too high (requires intense collision spike)."
             ],
             card_bg=COLOR_CARD_RED, accent_color=COLOR_RED)

    # -------------------------------------------------------------
    # SLIDE 8: 24-HOUR FINAL SPRINT ROADMAP
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Slide 08: 24-Hour Final Sprint Roadmap — Solving Choking Issues", "FINAL ROUND COMMITMENT")

    add_card(s8, Inches(0.8), Inches(1.8), Inches(2.2), Inches(5.0),
             "1. < 800ms Startup",
             "Cold Boot Fix",
             [
                 "Headless asset pre-warming.",
                 "Lazy bundle imports.",
                 "Eliminates render bottlenecks."
             ],
             accent_color=COLOR_CYAN)

    add_card(s8, Inches(3.18), Inches(1.8), Inches(2.2), Inches(5.0),
             "2. Live GPS Stream",
             "Continuous Stream",
             [
                 "WebSocket 5s updates.",
                 "Live speed & heading.",
                 "Polyline breadcrumbs."
             ],
             accent_color=COLOR_GREEN)

    add_card(s8, Inches(5.56), Inches(1.8), Inches(2.2), Inches(5.0),
             "3. Lock-Screen Hook",
             "Background Service",
             [
                 "Android Accessibility Background Service.",
                 "3-tap volume on locked OS screen."
             ],
             accent_color=COLOR_RED)

    add_card(s8, Inches(7.94), Inches(1.8), Inches(2.2), Inches(5.0),
             "4. Adaptive Fall Curve",
             "Kinematic Dual-Stage",
             [
                 "Free-fall weightlessness detection.",
                 "Secondary impact vector curve for slips."
             ],
             accent_color=COLOR_AMBER)

    add_card(s8, Inches(10.32), Inches(1.8), Inches(2.2), Inches(5.0),
             "5. Key Escrow",
             "Guardian Remote Access",
             [
                 "Asymmetric public/private key escrow.",
                 "Parents decrypt audio remotely."
             ],
             accent_color=COLOR_PURPLE)

    # -------------------------------------------------------------
    # SLIDE 9: COMPETITIVE ADVANTAGE
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Slide 09: Competitive Advantage — Why SafeGuard SOS Wins")

    rows = [
        ["Evaluation Metric", "Stock OS SOS (Apple/Google)", "Commercial Safety Apps", "SafeGuard SOS"],
        ["Zero-Install Web Tracking", "❌ Static SMS link only", "❌ Requires app install", "✅ Live Web Portal (Google Maps)"],
        ["Offline SMS Dispatch", "✅ Native SMS", "❌ Fails without 4G/5G", "✅ Direct SIM SMS + Cloud Sync"],
        ["Incident Evidence Capture", "❌ None", "❌ Manual photo only", "✅ Auto Ambient Audio + AES Vault"],
        ["Fall / Motion Detection", "⚠️ Watch / High-end only", "❌ None", "✅ Cross-Platform Sensor Engine"],
        ["Emergency Medical ID", "⚠️ Lock screen card", "❌ Paid subscription", "✅ Integrated Direct Broadcast"],
    ]

    table_shape = s9.shapes.add_table(len(rows), 4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    table = table_shape.table

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            cell.fill.solid()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            if r_idx == 0:
                cell.fill.fore_color.rgb = COLOR_CARD_BORDER
                p.font.bold = True
                p.font.size = Pt(13)
                p.font.color.rgb = COLOR_CYAN
            else:
                if c_idx == 3:
                    cell.fill.fore_color.rgb = COLOR_CARD_GREEN
                    p.font.bold = True
                    p.font.color.rgb = COLOR_GREEN
                elif r_idx % 2 == 1:
                    cell.fill.fore_color.rgb = COLOR_CARD_BG
                    p.font.color.rgb = COLOR_WHITE
                else:
                    cell.fill.fore_color.rgb = RGBColor(18, 22, 28)
                    p.font.color.rgb = COLOR_MUTED
                p.font.size = Pt(11)

    # -------------------------------------------------------------
    # SLIDE 10: MARKET OPPORTUNITY & SCALE
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "Slide 10: Market Opportunity & Scale — B2B and B2C Scalability")

    add_card(s10, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "1. Campus & Women Safety",
             "Direct B2C & B2B Institutional",
             [
                 "University campus security integration.",
                 "Rapid peer-to-peer buddy dispatch.",
                 "High demand for discreet, zero-panic trigger tools.",
                 "Freemium consumer tier for basic safety."
             ],
             accent_color=COLOR_CYAN)

    add_card(s10, Inches(4.85), Inches(1.8), Inches(3.6), Inches(5.0),
             "2. Lone Workers & Fleets",
             "Enterprise Workplace Compliance",
             [
                 "Delivery riders, night-shift personnel, and field technicians.",
                 "Automated crash & fall incident telemetry.",
                 "Centralized dispatcher web dashboard.",
                 "B2B SaaS subscription per seat."
             ],
             accent_color=COLOR_GREEN)

    add_card(s10, Inches(8.9), Inches(1.8), Inches(3.6), Inches(5.0),
             "3. Eldercare & Assisted Living",
             "Passive Health & Fall Monitoring",
             [
                 "Wearable BLE and phone accelerometer fall monitoring for elderly seniors.",
                 "Immediate family circle notifications with medical history payloads.",
                 "Peace of mind for caregivers."
             ],
             accent_color=COLOR_AMBER)

    # -------------------------------------------------------------
    # SLIDE 11: CONCLUSION
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)

    glow11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    glow11.fill.solid()
    glow11.fill.fore_color.rgb = COLOR_CARD_BG
    glow11.line.color.rgb = COLOR_GREEN
    glow11.line.width = Pt(2)

    tb11 = s11.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(11.0), Inches(5.5))
    tf11 = tb11.text_frame
    tf11.word_wrap = True

    p11_1 = tf11.paragraphs[0]
    p11_1.text = "🎯 HACKATHON ROUND 2 CONCLUSION"
    p11_1.font.size = Pt(13)
    p11_1.font.bold = True
    p11_1.font.color.rgb = COLOR_GREEN
    p11_1.space_after = Pt(8)

    p11_2 = tf11.add_paragraph()
    p11_2.text = "Engineered for Reality. Ready for the 24-Hour Final."
    p11_2.font.size = Pt(36)
    p11_2.font.bold = True
    p11_2.font.color.rgb = COLOR_WHITE
    p11_2.space_after = Pt(10)

    p11_3 = tf11.add_paragraph()
    p11_3.text = "We stress-tested the boundaries, identified our exact choke points, and established a concrete engineering plan to deliver an impenetrable production safety network."
    p11_3.font.size = Pt(14)
    p11_3.font.color.rgb = COLOR_MUTED
    p11_3.space_after = Pt(18)

    # TEAM SUMMARY FOOTER
    p11_team = tf11.add_paragraph()
    p11_team.text = "👥 TEAM: PIXEL PERFECT  |  TEAM ID: SM083"
    p11_team.font.size = Pt(14)
    p11_team.font.bold = True
    p11_team.font.color.rgb = COLOR_CYAN
    p11_team.space_after = Pt(6)

    p11_mem = tf11.add_paragraph()
    p11_mem.text = "Het Patel (240410149035)   •   Rudra Patel (240410149043)   •   Pritam Patel (240410149047)   •   Vedanshu Patel (240410149044)"
    p11_mem.font.size = Pt(12)
    p11_mem.font.color.rgb = COLOR_WHITE
    p11_mem.space_after = Pt(18)

    p11_4 = tf11.add_paragraph()
    p11_4.text = "🚀 Open for Questions & Live Sensor Demonstration"
    p11_4.font.size = Pt(18)
    p11_4.font.bold = True
    p11_4.font.color.rgb = COLOR_CYAN

    output_path = os.path.join(os.getcwd(), "SafeGuard_SOS_Round2_PitchDeck.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
